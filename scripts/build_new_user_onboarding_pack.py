"""Build the kid-level System3 onboarding CSV and PowerPoint deck.

The generated material is intentionally educational, not live-runtime proof.
It is grounded in the repository authorities named in ``SOURCES`` below.
"""

from __future__ import annotations

import csv
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "onboarding"
CSV_PATH = OUT / "SYSTEM3_NEW_USER_ALL_USE_CASES.csv"
PPTX_PATH = OUT / "SYSTEM3_NEW_USER_KID_LEVEL_GUIDE.pptx"

TODAY = date(2026, 8, 27).isoformat()
PROD_UI = "https://genesis-system3-web-doq2wplepa-el.a.run.app/ui/"
SAFETY = "ANALYZE_MODE=1; LIVE_TRADING_ENABLED=0; SYSTEM3_LIVE_TRADING_ALLOWED=0; AUTO_EXECUTE_TRADES=0"

SOURCES = [
    "AGENTS.md",
    "docs/control_plane/SYSTEM3_AGENT_RUNBOOK.md",
    "docs/authority/TEMPORAL_TRUTH_AND_LIVE_EVIDENCE_POLICY.md",
    "docs/authority/AUTONOMOUS_OPERATIONS_POLICY.md",
    "docs/project_control/SYSTEM3_MASTER_GOAL_LOCK.md",
    "docs/END_TO_END_ISSUES_SOLUTIONS_AGENT_POLICY.md",
    "docs/CONTINUOUS_CLOSURE_SYSTEM.md",
    "docs/PREFLIGHT_CONTROL_PLANE.md",
    "docs/project_control/REPO_CLEAN_FORENSIC_TOOLKIT.md",
    "docs/project_control/PREDICTION_WORLD_CLASS_BENCHMARK_POLICY.md",
    "dashboard/frontend/src/components/Sidebar.tsx",
    "dashboard/frontend/src/App.tsx",
    ".env.example",
]

COLUMNS = [
    "use_case_id", "audience", "journey_stage", "category", "use_case",
    "kid_level_goal", "when_to_use", "simple_steps", "ui_tab_or_surface",
    "inputs_needed", "expected_output", "proof_or_acceptance", "status_words",
    "safety_boundary", "command_or_url", "authority_source", "user_action",
]


def row(i, audience, stage, category, use_case, goal, when, steps, surface,
        inputs, output, proof, statuses="PASS | PARTIAL | FAIL | BLOCKED",
        safety="Read-only or PAPER/ANALYZE only; never place a real order.",
        command="", source="AGENTS.md", user_action="None for routine safe work"):
    return {
        "use_case_id": f"UC-{i:03d}", "audience": audience,
        "journey_stage": stage, "category": category, "use_case": use_case,
        "kid_level_goal": goal, "when_to_use": when, "simple_steps": steps,
        "ui_tab_or_surface": surface, "inputs_needed": inputs,
        "expected_output": output, "proof_or_acceptance": proof,
        "status_words": statuses, "safety_boundary": safety,
        "command_or_url": command, "authority_source": source,
        "user_action": user_action,
    }


def build_rows():
    r = []
    add = lambda *a, **k: r.append(row(len(r) + 1, *a, **k))
    add("Everyone", "Start", "Orientation", "Understand System3", "See it as a careful market-science robot, not a money machine.", "Your first five minutes.", "Read the safety card > learn the four layers > open the guide.", "Guide", "None", "A correct mental model", "Can explain DATA > THINK > PAPER > PROVE", source="docs/SYSTEM3_CORE_TRADING_GOAL_AND_ARCHITECTURE.md")
    add("Everyone", "Start", "Safety", "Confirm the four safety locks", "Make sure the robot can watch and practice but cannot trade real money.", "Before every run or proof.", "Check all four values are exactly as listed.", "Environment / Gates", "Environment values", "All locks OFF except ANALYZE_MODE=1", "Exact values match", statuses="SAFE | UNSAFE", safety="Stop immediately if a live/order flag is enabled.", command=SAFETY, source="AGENTS.md")
    add("Viewer", "Start", "Access", "Open the production dashboard", "Open the real cloud screen.", "To view current production.", "Start a fresh browser session > open URL > note the capture time.", "Production UI", "Internet + browser", "Dashboard shell", "Fresh page opened after request start; serving SHA checked for current claims", command=PROD_UI, source="docs/authority/TEMPORAL_TRUTH_AND_LIVE_EVIDENCE_POLICY.md")
    add("Viewer", "Start", "Truth", "Tell live truth from old proof", "A photo from yesterday is not what is happening now.", "Whenever words like now/current/live appear.", "Record UTC start > create new browser session > capture UI + API > compare.", "All surfaces", "Fresh timestamps", "Time-bounded evidence", "Evidence age and capture time reported", statuses="FRESH | STALE | HISTORICAL | UNVERIFIED", source="docs/authority/TEMPORAL_TRUTH_AND_LIVE_EVIDENCE_POLICY.md")
    add("New local user", "Setup", "Repository", "Open the correct repo", "Stand in the right workshop.", "Before commands or edits.", "Open PowerShell > change to this repo > run git status.", "Repository", "Local clone", "Correct root and visible changes", "No unrelated changes overwritten", command=r"Set-Location C:\Users\ADMIN\Genesis_System3\Genesis_System3; git status --short", source="AGENTS.md")
    add("New local user", "Setup", "Dependencies", "Check tools", "Make sure the toolbox exists.", "First local setup or after machine changes.", "Check Python > Node > npm > virtual environment.", "Terminal", "Installed tools", "Versions print without error", "Python and Node checks succeed", command="python --version; node --version; npm --version", source="dashboard/README.md")
    add("New local user", "Run", "Backend", "Start local API brain", "Wake up the part that answers questions.", "Local development.", "Set safety flags > start uvicorn > keep window open.", "Local API", "Python environment", "API listens on 127.0.0.1:8000", "GET /api/health returns 200", command=r".\.venv\Scripts\python.exe -m uvicorn dashboard.backend.app:app --host 127.0.0.1 --port 8000", source="README.md")
    add("New local user", "Run", "Dashboard", "Start/open local UI", "Open the control-room screen.", "After backend starts.", "Use the repo dashboard launcher or frontend dev server > open the printed URL.", "Local UI", "Backend + frontend dependencies", "Dashboard appears", "No red browser errors; LIVE OFF visible", command=r".\scripts\run_dashboard.ps1", source="dashboard/README.md")
    add("Viewer", "Use", "Navigation", "Use a tab link", "Jump straight to one room.", "Sharing a specific view.", "Add ?tab=<tab-id> to the production UI URL.", "Any canonical tab", "Tab ID", "Requested tab opens", "URL and selected tab agree", command=PROD_UI + "?tab=decision-intel", source="dashboard/frontend/src/App.tsx")

    tabs = [
        ("decision-intel", "Decision Intel", "Ask: what is the system thinking, and why?", "Current decision evidence and reasons"),
        ("truth", "Truth", "Check whether labels match real sources and freshness.", "Truth/freshness summary"),
        ("genesis", "Genesis Brain", "See model evidence without pretending missing data exists.", "Model status, confidence, reasons or WAITING"),
        ("e2e-proof", "E2E Proof", "Follow proof from code to cloud to browser.", "Proof-chain status"),
        ("overview", "Overview", "Get the big picture first.", "Mode, market, health and summary"),
        ("sim-live", "Live Simulation", "Watch a safe simulation, not real execution.", "Simulation state and events"),
        ("options-intel", "Options Intel", "Compare option opportunities and evidence.", "Options intelligence summary"),
        ("chain", "Option Chain", "Read strikes, CE/PE, OI, volume and freshness.", "Chain rows or an honest unavailable state"),
        ("signals", "Signals", "See suggestions plus the reason and block gates.", "Signal/no-trade decision"),
        ("trade", "Trade Plan", "Review a plan without sending a broker order.", "Read-only plan, entry/stop/target evidence"),
        ("paper", "Paper Trading", "Practice trades with pretend money.", "Paper orders, ledger and P&L"),
        ("positions", "Positions", "Separate Dhan account truth from paper positions.", "Clearly sourced position rows"),
        ("risk-scenarios", "Risk & Scenarios", "Ask what could go wrong before acting.", "Scenario and risk evidence"),
        ("multibagger", "Multibagger Research", "Study long-horizon candidates as research only.", "Candidate evidence or honest empty state"),
        ("prediction-audit", "Prediction Audit", "Check whether predictions were made before outcomes.", "Prediction lineage and evaluation"),
        ("performance", "Performance", "Measure results with costs and drawdown.", "Performance metrics with provenance"),
        ("ml", "ML", "Inspect models; do not trust one accuracy number.", "Model quality, drift and calibration evidence"),
        ("data-integrity", "Data Integrity", "Check if data is fresh, real and complete.", "Source, age, gaps and quality"),
        ("broker", "Broker", "Check Dhan connection in read-only mode.", "Auth/connection/funds/holdings status"),
        ("alerts", "Alerts", "See what needs attention without panic.", "Deduplicated warnings and severity"),
        ("system", "System", "Inspect services, versions and diagnostics.", "Health and component status"),
        ("gates", "Gates", "See which locks block readiness.", "Gate pass/fail and blockers"),
    ]
    for tab_id, label, goal, output in tabs:
        add("Viewer / Paper operator", "Use", "Dashboard tab", f"Use {label}", goal, f"When the {label} question matters.", f"Open {label} > read source/as-of/status > inspect warnings > do not turn blanks into guesses.", tab_id, "Fresh dashboard session", output, "UI meaning matches same-session API; empty/stale is explicit", command=f"{PROD_UI}?tab={tab_id}", source="dashboard/frontend/src/components/Sidebar.tsx")

    add("Paper operator", "Daily", "Market", "Before-market check", "Check the playground before practice starts.", "Before Indian market hours.", "Safety locks > health > Dhan read-only > data freshness > chain > gates.", "Overview / Broker / Data Integrity / Gates", "Fresh session", "Ready, partial, or blocked checklist", "No hidden red/stale state", source="AGENTS.md")
    add("Paper operator", "Daily", "Market", "During-market watch", "Watch data move without touching real money.", "During market session.", "Refresh health > compare UI/API > watch stale/429/401/5xx > keep paper separate.", "Overview / Chain / Signals / Paper", "Market data", "Timestamped observations", "Freshness and source stay explicit", source="docs/control_plane/SYSTEM3_AGENT_RUNBOOK.md")
    add("Paper operator", "Daily", "Reconciliation", "End-of-day review", "Compare guesses with what actually happened.", "After the trading day.", "Join predictions to outcomes > review selected/rejected > include costs > classify misses.", "Prediction Audit / Performance", "Prediction and paper ledgers", "PASS/PARTIAL/FAIL/BLOCKED verdict", "Row-level lineage exists", source="docs/SYSTEM3_CORE_TRADING_GOAL_AND_ARCHITECTURE.md")
    add("Paper operator", "Use", "Paper", "Create a paper-only learning loop", "Practice, measure, improve.", "Testing a signal or strategy.", "Record prediction > simulate realistic fill > apply risk > exit > reconcile.", "Signals / Paper / Performance", "Fresh quote + strategy", "Auditable paper trade", "No broker mutation; costs/slippage recorded", safety="Paper only. A paper button or plan must never call a live order route.", source="docs/project_control/PREDICTION_WORLD_CLASS_BENCHMARK_POLICY.md")
    add("Analyst", "Research", "Options", "Compare NIFTY-family chains", "Check more than one option playground.", "Options/chain research.", "Check NIFTY > BANKNIFTY > FINNIFTY > MIDCPNIFTY separately.", "Options Intel / Chain", "Fresh contracts and expiries", "Per-underlying evidence", "Symbol, source, expiry, strikes and timestamps verified", source="docs/control_plane/SYSTEM3_AGENT_RUNBOOK.md")
    add("Analyst", "Research", "Prediction", "Benchmark a model challenger", "Let a new brain compete fairly.", "Any material feature/model change.", "Fix data truth > prevent leakage > same-window tournament > costs > calibration > robustness.", "ML / Prediction Audit", "Point-in-time dataset", "Champion/challenger matrix", "All promotion gates green; simple baselines included", safety="No automatic LIVE enablement or capital deployment.", source="docs/project_control/PREDICTION_WORLD_CLASS_BENCHMARK_POLICY.md")
    add("Analyst", "Research", "Data", "Add a feature safely", "Teach one new clue and prove it helps.", "Feature engineering.", "Record source/version/as-of/availability > leakage test > missingness > OOS ablation.", "Data Integrity / ML", "Point-in-time feature", "Versioned feature row", "Incremental OOS value proven", source="docs/project_control/PREDICTION_WORLD_CLASS_BENCHMARK_POLICY.md")
    add("Developer", "Change", "Code", "Fix a bug", "Find the broken brick and replace only that brick.", "Verified repo or runtime issue.", "Reproduce > root cause > regression test > smallest fix > focused tests > PR.", "Repository / PR", "Current main + issue evidence", "Reviewed code change", "Exact-head required CI green", source="docs/control_plane/SYSTEM3_AGENT_RUNBOOK.md")
    add("Developer", "Verify", "Tests", "Run focused tests", "Check the part you changed.", "Before broad CI.", "Choose nearby unit/contract tests > run > read the first real failure.", "Terminal", "Changed files", "Test results", "Relevant tests pass", command=r".\.venv\Scripts\python.exe -m pytest <test_path> -q", source="docs/control_plane/SYSTEM3_AGENT_RUNBOOK.md")
    add("Developer", "Verify", "Frontend", "Build and browser-test UI", "Make sure the screen still fits and works.", "Frontend changes.", "Build TypeScript/Vite > run Playwright > inspect screenshots/console/network.", "Frontend", "Node dependencies", "Build + browser evidence", "No regressions at applicable viewports", command="npm run build; npm run test:dashboard", source="package.json")
    add("Maintainer", "Control", "Preflight", "Generate current control-plane snapshot", "Check traffic lights before moving production.", "Before merge/deploy/production transition.", "Run fresh preflight > inspect main/PR/workflows/artifacts/issues.", "Control plane", "GitHub/GCP access", "Fresh snapshot", "Critical claims independently revalidated", command=r".\.venv\Scripts\python.exe scripts\system3_preflight_control_plane.py", source="docs/PREFLIGHT_CONTROL_PLANE.md")
    add("Maintainer", "Control", "CI", "Inspect a failed workflow", "Open the red box and read the label inside.", "Current-main or active-PR failure.", "Open run > failed job > failed step > log/artifact > fix cause.", "GitHub Actions", "Run ID", "Root-cause evidence", "No diagnosis from red icon alone", command="gh run view <run-id> --log-failed", source="docs/control_plane/SYSTEM3_AGENT_RUNBOOK.md")
    add("Maintainer", "Release", "Merge", "Merge an exact proven head", "Only join code that passed its own test.", "Required CI green on exact PR head.", "Re-read runbook > fetch > verify SHA/checks > merge.", "GitHub PR", "Exact head SHA", "Merge SHA", "Mandatory checks green for exact head", source="docs/control_plane/SYSTEM3_AGENT_RUNBOOK.md")
    add("Maintainer", "Release", "Deploy", "Verify Cloud Run deployment", "Make sure cloud serves the new box.", "Immediately after merge.", "Check project/region/service > revision/traffic > serving SHA > safety env.", "GCP Cloud Run", "Merge SHA", "Serving revision proof", "100% intended traffic and exact provenance", safety="No service-account JSON keys; WIF only.", source="docs/control_plane/SYSTEM3_AGENT_RUNBOOK.md")
    add("Maintainer", "Release", "Production proof", "Capture fresh 22-tab proof", "Take a new report card after the fix reaches cloud.", "After serving SHA converges.", "New browser > actual GCP URL > all 22 tabs > APIs > console/network > end SHA.", "Production UI", "Deployed SHA", "Timestamped semantic proof pack", "UI/API agree and evidence is request-scoped", command=r".\.venv\Scripts\python.exe scripts\gcp_live_ui_snapshot.py", source="docs/authority/TEMPORAL_TRUTH_AND_LIVE_EVIDENCE_POLICY.md")
    add("Maintainer", "Operate", "Broker", "Investigate Dhan authentication", "Check the key works without showing the key.", "Broker is disconnected or 401-like.", "Use bounded read-only status > inspect redacted logs > dedicated recovery authority only if needed.", "Broker", "Authorized runtime access", "AUTH_OK/DEGRADED/BLOCKED reason", "No secret payload exposed", safety="Never paste tokens; rotation only through bounded Dhan authority.", source="AGENTS.md")
    add("Maintainer", "Operate", "Scheduler", "Check recurring jobs", "Make sure the alarm clock rings and does useful work.", "Missing/stale forecasts, ranks, or signals.", "Check scheduler transport > job date > output freshness > business result.", "System / Alerts", "Fresh scheduler/log evidence", "Transport and business-health verdicts", "Healthy transport is not confused with correct business date", source="reports/coordination/COMMAND_CENTER.md")
    add("Maintainer", "Operate", "Repo cleanup", "Find safe storage savings", "Sort trash only after proving it is trash.", "Repo size/duplicate/storage request.", "Run canonical toolkit > read current executive decision > act only on DELETE_PROVEN_100 rows via PR.", "Repository", "Fresh current-main scan", "Report-only cleanup decision", "Storage layers separated; no history rewrite", safety="Toolkit never deletes source, artifacts, history, or secrets.", command=r".\.venv\Scripts\python.exe scripts\system3_repo_clean_forensic_toolkit.py", source="docs/project_control/REPO_CLEAN_FORENSIC_TOOLKIT.md")
    add("Everyone", "Troubleshoot", "Statuses", "Read status words correctly", "Know whether to smile, wait, or investigate.", "Any report or dashboard card.", "PASS=proven > PARTIAL=some proof > FAIL=wrong > BLOCKED=missing dependency > WAITING=in progress.", "All surfaces", "Status + evidence", "Correct next action", "Status has timestamp and named evidence", statuses="PASS | PARTIAL | FAIL | BLOCKED | WAITING | DEGRADED | STALE", source="docs/control_plane/SYSTEM3_AGENT_RUNBOOK.md")
    add("Everyone", "Troubleshoot", "HTTP", "Understand common web errors", "Translate the number into a next check.", "API/UI failure.", "401=auth > 404=route > 429=rate limit > 5xx=server > timeout=network/service.", "Browser / API", "Status code + response", "Classified problem", "Read response/log; do not guess", statuses="200 | 401 | 404 | 429 | 5xx | TIMEOUT", source="docs/control_plane/SYSTEM3_AGENT_RUNBOOK.md")
    add("Everyone", "Troubleshoot", "Empty data", "Treat blank as a clue, not zero.", "Find out why the box is empty before deciding what it means.", "Table/chart has no rows.", "Check source > timestamp > market state > API > error banner > lineage.", "Any data tab", "UI + API response", "EMPTY/WAITING/STALE reason", "No synthetic fill presented as live", statuses="EMPTY | WAITING | STALE | ERROR", source="AGENTS.md")
    add("Everyone", "Security", "Secrets", "Keep passwords and broker keys private", "Never show the secret sticker.", "Logs, chats, screenshots, commits.", "Use templates/secret manager > redact > scan > rotate only through authority.", "All surfaces", "Secret names, never values", "No exposed secret", "Secret scan and redacted evidence", safety="Never commit or paste Dhan/GitHub/GCP secret payloads.", source="AGENTS.md")
    add("Everyone", "Safety", "Forbidden", "Know the red line", "Never let a learning task become a real-money action.", "Always.", "Do not place/modify/cancel/square-off real orders; do not enable LIVE.", "All surfaces", "None", "Action refused and safely redirected", "Safety flags unchanged", statuses="FORBIDDEN", safety="Human break-glass is required for LIVE or real orders.", source="AGENTS.md", user_action="Only explicit human break-glass/account-level approval")
    add("Agent", "Communicate", "Progress", "Write a useful work update", "Tell the next helper exactly where the work is.", "Every active remediation update.", "State STATUS > IN_PROGRESS > CURRENT_STEP > NEXT_ACTION > USER_ACTION.", "Issue / PR / report", "Fresh evidence", "Five-field update", "Claims are timestamped and scoped", source="AGENTS.md")
    add("Agent", "Close", "Completion", "Close only after user-visible proof", "Do not call the toy fixed while it is still broken on screen.", "After code/CI/deploy.", "Re-read runbook > prove serving SHA > new UI/API proof > ledger > close.", "Completion ledger", "Fresh production evidence", "Evidence-backed closure", "No required work remains", source="docs/control_plane/SYSTEM3_AGENT_RUNBOOK.md")
    add("Everyone", "Reference", "CSV", "Use this CSV as a checklist", "Filter the big map into your small job.", "Training, SOPs, QA, audits, handoffs.", "Filter by audience/category/tab/status > follow steps > record proof.", "This CSV", "Spreadsheet app", "Personalized checklist", "Original row IDs retained", source="docs/onboarding/SYSTEM3_NEW_USER_ALL_USE_CASES.csv")
    return r


NAVY = RGBColor(7, 18, 31)
BLUE = RGBColor(59, 140, 255)
CYAN = RGBColor(32, 211, 238)
GREEN = RGBColor(24, 215, 130)
AMBER = RGBColor(245, 165, 36)
RED = RGBColor(255, 73, 100)
WHITE = RGBColor(243, 247, 252)
MUTED = RGBColor(167, 184, 204)
CARD = RGBColor(15, 31, 49)


def rect(slide, x, y, w, h, fill, radius=True, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line or fill
    return shp


def txt(slide, text, x, y, w, h, size=18, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text; run.font.name = "Aptos"; run.font.size = Pt(size)
    run.font.bold = bold; run.font.color.rgb = color
    return box


def base_slide(prs, title, kicker="SYSTEM3 • NEW USER GUIDE"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = NAVY
    rect(s, 0, 0, 13.333, .08, BLUE, False)
    txt(s, kicker, .55, .30, 6.5, .25, 9, CYAN, True)
    txt(s, title, .55, .62, 12.1, .62, 26, WHITE, True)
    txt(s, f"Educational guide • generated {TODAY} • not proof of current runtime state", .55, 7.15, 8.8, .2, 8, MUTED)
    txt(s, "LIVE OFF • PAPER / ANALYZE", 10.15, 7.08, 2.6, .28, 9, GREEN, True, PP_ALIGN.RIGHT)
    return s


def card(slide, x, y, w, h, title, body, accent=BLUE, icon=None, body_size=14):
    rect(slide, x, y, w, h, CARD, True, RGBColor(30, 56, 82))
    rect(slide, x, y, .07, h, accent, False)
    if icon:
        txt(slide, icon, x + .18, y + .18, .42, .35, 18, accent, True, PP_ALIGN.CENTER)
        tx = x + .68
    else:
        tx = x + .25
    txt(slide, title, tx, y + .18, w - (tx-x) - .18, .35, 15, WHITE, True)
    txt(slide, body, x + .25, y + .68, w - .48, h - .82, body_size, MUTED)


def bullets(slide, items, x, y, w, h, size=15, color=MUTED, accent=CYAN):
    each = h / max(len(items), 1)
    for i, item in enumerate(items):
        txt(slide, "●", x, y + i*each, .22, each, max(8, size-4), accent, True)
        txt(slide, item, x + .28, y + i*each, w - .28, each, size, color)


def build_pptx():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
    rect(s, 0, 0, 13.333, .10, BLUE, False)
    txt(s, "GENESIS SYSTEM3", .65, .65, 4.5, .35, 13, CYAN, True)
    txt(s, "The kid-simple\nnew user guide", .65, 1.22, 7.4, 1.55, 34, WHITE, True)
    txt(s, "See data. Think carefully. Practice on paper. Prove what happened.", .68, 3.03, 6.8, .75, 18, MUTED)
    card(s, 8.35, .85, 4.2, 1.18, "1 • LOOK", "Fresh market and system data", BLUE, "◉")
    card(s, 8.35, 2.20, 4.2, 1.18, "2 • THINK", "Signals, models and risk", CYAN, "◆")
    card(s, 8.35, 3.55, 4.2, 1.18, "3 • PRACTICE", "Paper trades only", AMBER, "✎")
    card(s, 8.35, 4.90, 4.2, 1.18, "4 • PROVE", "UI + API + time + lineage", GREEN, "✓")
    txt(s, "For viewers • paper operators • analysts • developers • maintainers", .68, 6.55, 7.0, .3, 12, WHITE, True)
    txt(s, f"{TODAY} • Educational, not live evidence", .68, 6.95, 5.5, .2, 9, MUTED)

    s = base_slide(prs, "The one-sentence idea")
    txt(s, "System3 is a careful market-science robot.", .72, 1.55, 11.8, .65, 28, WHITE, True, PP_ALIGN.CENTER)
    txt(s, "It watches real information, makes explainable guesses, practices with pretend money, and checks whether those guesses were right.", 1.25, 2.45, 10.8, 1.0, 20, MUTED, False, PP_ALIGN.CENTER)
    card(s, 1.05, 4.10, 3.35, 1.25, "It IS", "Analyzer • paper lab • proof system", GREEN, "✓")
    card(s, 4.95, 4.10, 3.35, 1.25, "It IS NOT", "A promise of profit or a magic crystal ball", AMBER, "!")
    card(s, 8.85, 4.10, 3.35, 1.25, "It MUST NOT", "Place, change, cancel, or close real orders", RED, "×")

    s = base_slide(prs, "The four safety locks")
    locks = [("ANALYZE_MODE", "1", GREEN), ("LIVE_TRADING_ENABLED", "0", BLUE), ("SYSTEM3_LIVE_TRADING_ALLOWED", "0", BLUE), ("AUTO_EXECUTE_TRADES", "0", BLUE)]
    for i, (name, val, col) in enumerate(locks):
        x = .72 + (i % 2)*6.15; y = 1.45 + (i//2)*2.0
        rect(s, x, y, 5.72, 1.55, CARD, True, col)
        txt(s, "LOCKED", x+.25, y+.22, 1.15, .25, 10, col, True)
        txt(s, name, x+.25, y+.58, 4.2, .32, 15, WHITE, True)
        txt(s, f"= {val}", x+4.55, y+.48, .75, .55, 26, col, True, PP_ALIGN.RIGHT)
    txt(s, "If any LIVE/order lock changes, stop. Real-money actions need explicit human break-glass authority.", .9, 5.72, 11.5, .55, 17, RED, True, PP_ALIGN.CENTER)

    s = base_slide(prs, "The truth ladder: what counts as ‘now’?")
    ladder = [("1", "Fresh production browser", "New session after the request", GREEN), ("2", "Same-session APIs", "Values, source, time and errors", CYAN), ("3", "Runtime logs + metadata", "Same observation window", BLUE), ("4", "Serving revision / SHA", "Exact deployed identity", AMBER), ("5", "Source code", "What should happen", MUTED), ("6", "Stored reports", "History until freshly revalidated", RED)]
    for i,(n,t,b,c) in enumerate(ladder):
        y=1.35+i*.82; rect(s,.85,y,11.65,.65,CARD,True,RGBColor(30,56,82)); rect(s,.85,y,.62,.65,c,True,c)
        txt(s,n,.85,y+.10,.62,.28,14,NAVY,True,PP_ALIGN.CENTER); txt(s,t,1.68,y+.12,3.9,.25,14,WHITE,True); txt(s,b,5.35,y+.12,6.75,.25,13,MUTED)

    s = base_slide(prs, "Who are you today?")
    roles=[("VIEWER","Read cards, timestamps, sources and warnings","No changes",BLUE), ("PAPER OPERATOR","Run safe practice and daily checks","Paper only",AMBER), ("ANALYST","Study data, signals, risks and models","Research only",CYAN), ("DEVELOPER","Fix code and tests on a branch","No production mutation by default",GREEN), ("MAINTAINER","CI, merge, deploy and fresh proof","Only after gates",RGBColor(168,85,247))]
    for i,(t,b,tag,c) in enumerate(roles):
        x=.65+(i%3)*4.22; y=1.45+(i//3)*2.25; card(s,x,y,3.82,1.82,t,b+"\n\nBoundary: "+tag,c,str(i+1),13)

    s = base_slide(prs, "Your first 10 minutes")
    steps=[("1", "Open the correct repo", r"C:\Users\ADMIN\Genesis_System3\Genesis_System3"), ("2", "Confirm the four locks", "1 / 0 / 0 / 0"), ("3", "Open the dashboard", "Production for viewing; local for development"), ("4", "Read Overview + Truth + Gates", "Big picture, evidence, blockers"), ("5", "Say what you know", "Include capture time and evidence age")]
    for i,(n,t,b) in enumerate(steps):
        y=1.38+i*.98; rect(s,.95,y,11.45,.78,CARD,True,RGBColor(30,56,82)); rect(s,.95,y,.78,.78,BLUE,True,BLUE); txt(s,n,.95,y+.18,.78,.30,18,WHITE,True,PP_ALIGN.CENTER); txt(s,t,1.95,y+.13,4.2,.28,15,WHITE,True); txt(s,b,6.15,y+.13,5.85,.4,13,MUTED)

    s = base_slide(prs, "The dashboard: 22 rooms")
    groups=[("START HERE","Decision Intel • Truth • Genesis • E2E Proof • Overview • Live Simulation",BLUE), ("MARKET","Options Intel • Option Chain • Signals • Trade Plan • Paper Trading • Positions",CYAN), ("ANALYSIS","Risk & Scenarios • Multibagger Research • Prediction Audit • Performance • ML",AMBER), ("SYSTEM","Data Integrity • Broker • Alerts • System • Gates",GREEN)]
    for i,(t,b,c) in enumerate(groups):
        x=.72+(i%2)*6.15; y=1.45+(i//2)*2.45; card(s,x,y,5.72,2.05,t,b,c,str(i+1),15)

    s = base_slide(prs, "Start-here rooms")
    items=[("Decision Intel","What is the system thinking—and why?"), ("Truth","Are labels, sources and freshness honest?"), ("Genesis Brain","What model evidence exists or is WAITING?"), ("E2E Proof","Can we trace code → cloud → browser?"), ("Overview","What is the big picture?"), ("Live Simulation","What happens in safe simulation?")]
    for i,(t,b) in enumerate(items): card(s,.65+(i%3)*4.22,1.4+(i//3)*2.35,3.82,1.85,t,b,BLUE,str(i+1),13)

    s = base_slide(prs, "Market rooms")
    items=[("Options Intel","Compare opportunity evidence"), ("Option Chain","Strikes, CE/PE, OI, volume, age"), ("Signals","Suggestion + reasons + blockers"), ("Trade Plan","Plan only—no broker order"), ("Paper Trading","Pretend fills, ledger and P&L"), ("Positions","Dhan truth kept separate from paper")]
    for i,(t,b) in enumerate(items): card(s,.65+(i%3)*4.22,1.4+(i//3)*2.35,3.82,1.85,t,b,CYAN,str(i+1),13)

    s = base_slide(prs, "Analysis + system rooms")
    left=["Risk & Scenarios — what could go wrong?","Multibagger — long-horizon research","Prediction Audit — guess before result?","Performance — costs, drawdown, expectancy","ML — drift, calibration, evidence"]
    right=["Data Integrity — real, fresh, complete?","Broker — Dhan read-only status","Alerts — deduplicated warnings","System — services and versions","Gates — which locks block readiness?"]
    card(s,.7,1.42,5.9,4.95,"ANALYSIS", "",AMBER,"A")
    bullets(s,left,1.05,2.15,5.15,3.55,14,WHITE,AMBER)
    card(s,6.75,1.42,5.9,4.95,"SYSTEM", "",GREEN,"S")
    bullets(s,right,7.1,2.15,5.15,3.55,14,WHITE,GREEN)

    s = base_slide(prs, "One safe daily loop")
    flow=[("BEFORE", "Locks → health → broker → data → gates", BLUE), ("DURING", "Watch freshness, errors, chain, signals and paper", CYAN), ("AFTER", "Predictions → outcomes → costs → misses → verdict", AMBER), ("LEARN", "Small change → fair test → evidence → keep/rollback", GREEN)]
    for i,(t,b,c) in enumerate(flow):
        x=.7+i*3.15; rect(s,x,2.0,2.75,2.4,CARD,True,c); txt(s,str(i+1),x+.22,2.25,.48,.5,24,c,True); txt(s,t,x+.82,2.28,1.65,.32,15,WHITE,True); txt(s,b,x+.25,3.05,2.25,.95,13,MUTED,False,PP_ALIGN.CENTER)
        if i<3: txt(s,"→",x+2.77,2.85,.38,.4,22,WHITE,True,PP_ALIGN.CENTER)
    txt(s,"Every observation needs source + as-of time + freshness + status.",1.2,5.3,10.9,.55,18,WHITE,True,PP_ALIGN.CENTER)

    s = base_slide(prs, "How to read a status")
    statuses=[("PASS","The named check is proven",GREEN), ("PARTIAL","Some proof works; gaps remain",AMBER), ("FAIL","The check ran and found wrongness",RED), ("BLOCKED","A required dependency is missing",RGBColor(168,85,247)), ("WAITING","Work/state is not finished",BLUE), ("STALE","Evidence is too old for ‘now’",MUTED)]
    for i,(t,b,c) in enumerate(statuses): card(s,.65+(i%3)*4.22,1.4+(i//3)*2.35,3.82,1.85,t,b,c,"●",13)

    s = base_slide(prs, "When a box is empty")
    txt(s,"EMPTY ≠ ZERO ≠ BROKEN",1.0,1.45,11.3,.55,27,WHITE,True,PP_ALIGN.CENTER)
    checks=[("SOURCE","Which API/file/provider?"), ("TIME","When was it observed?"), ("MARKET","Open, closed, holiday?"), ("ERROR","401, 404, 429, 5xx?"), ("LINEAGE","Can the row be traced?"), ("LABEL","EMPTY, WAITING, STALE, ERROR?")]
    for i,(t,b) in enumerate(checks): card(s,.65+(i%3)*4.22,2.35+(i//3)*1.72,3.82,1.3,t,b,AMBER,str(i+1),12)
    txt(s,"Never invent synthetic values and call them live.",1.3,6.15,10.7,.4,18,RED,True,PP_ALIGN.CENTER)

    s = base_slide(prs, "Prediction: a fair science contest")
    chain=["DATA","LINEAGE","FEATURE","LEAKAGE","LABEL","TOURNAMENT","OOS","COST","CALIBRATION","ROBUSTNESS","PAPER","UI/API","CHAMPION"]
    for i,t in enumerate(chain):
        col=i%7; row_=i//7; x=.55+col*1.8; y=1.6+row_*1.55
        rect(s,x,y,1.55,.82,CARD,True, GREEN if i==12 else BLUE)
        txt(s,t,x+.08,y+.23,1.39,.25,10,WHITE,True,PP_ALIGN.CENTER)
        if i<12 and col<6: txt(s,"→",x+1.55,y+.22,.25,.25,12,MUTED,True,PP_ALIGN.CENTER)
    bullets(s,["Simple baselines must compete.","One accuracy number is never enough.","Use realistic costs, uncertainty and regime splits.","Self-learning means bounded retraining—not self-modifying LIVE code."],1.0,4.85,11.3,1.25,14,WHITE,CYAN)

    s = base_slide(prs, "Developer path: smallest safe change")
    steps=["Verify fresh symptom","Find root cause","Add regression test","Make smallest fix","Run focused checks","Exact-head CI","Merge","Deploy","New live proof","Close ledger"]
    for i,t in enumerate(steps):
        x=.6+(i%5)*2.5; y=1.55+(i//5)*2.05
        rect(s,x,y,2.15,1.3,CARD,True,BLUE if i<5 else GREEN)
        txt(s,f"{i+1:02d}",x+.15,y+.15,.48,.3,13,CYAN,True)
        txt(s,t,x+.18,y+.55,1.78,.5,13,WHITE,True,PP_ALIGN.CENTER)
    txt(s,"Green CI proves its tests—not automatically the user-visible result.",1.0,5.95,11.3,.42,17,AMBER,True,PP_ALIGN.CENTER)

    s = base_slide(prs, "Production proof: the same-SHA chain")
    nodes=[("GITHUB","Exact head + checks"), ("BUILD","Immutable artifact"), ("CLOUD RUN","Revision + traffic"), ("URL","Serving SHA"), ("BROWSER","22 tabs + F12"), ("SEMANTICS","UI = API truth")]
    for i,(t,b) in enumerate(nodes):
        x=.48+i*2.12; rect(s,x,2.0,1.82,2.2,CARD,True,GREEN if i==5 else BLUE)
        txt(s,t,x+.1,2.35,1.62,.28,13,WHITE,True,PP_ALIGN.CENTER); txt(s,b,x+.15,3.05,1.52,.6,11,MUTED,False,PP_ALIGN.CENTER)
        if i<5: txt(s,"→",x+1.82,2.82,.3,.35,18,CYAN,True,PP_ALIGN.CENTER)
    txt(s,"A pre-fix screenshot becomes history after the fix. Capture again.",1.0,5.15,11.3,.5,18,AMBER,True,PP_ALIGN.CENTER)

    s = base_slide(prs, "Troubleshooting without guessing")
    errs=[("401","Authentication/session"), ("404","Route not present"), ("429","Rate limit/backoff"), ("5xx","Server failure"), ("TIMEOUT","Network/service delay"), ("STALE","Old data/evidence")]
    for i,(t,b) in enumerate(errs): card(s,.65+(i%3)*4.22,1.4+(i//3)*2.35,3.82,1.85,t,b,RED if t in ("5xx","401") else AMBER,"!",13)

    s = base_slide(prs, "Secrets and red lines")
    card(s,.72,1.5,5.72,4.65,"DO", "Use .env.example as a template\nUse approved secret stores\nRedact logs and screenshots\nUse Dhan read-only checks\nUse WIF for GitHub → GCP",GREEN,"✓",15)
    card(s,6.88,1.5,5.72,4.65,"DO NOT", "Paste broker tokens in chat\nCommit credentials\nCreate service-account JSON keys\nEnable LIVE or auto-execute\nPlace, change, cancel or square-off real orders",RED,"×",15)

    s = base_slide(prs, "Use the CSV as your master checklist")
    card(s,.72,1.45,4.0,4.95,"FILTER", "Audience\nJourney stage\nCategory\nDashboard tab\nStatus word",BLUE,"1",16)
    card(s,4.95,1.45,4.0,4.95,"FOLLOW", "Kid-level goal\nWhen to use\nSimple steps\nInputs and outputs\nSafety boundary",CYAN,"2",16)
    card(s,9.18,1.45,3.45,4.95,"PROVE", "Acceptance check\nAuthority source\nUser action\nKeep the UC ID",GREEN,"3",16)

    s = base_slide(prs, "The tiny cheat sheet")
    bullets(s,["Start with Overview, Truth and Gates.","For ‘now’, create fresh browser + API evidence.","For market work, stay PAPER / ANALYZE only.","For blank data, check source, time, market and errors.","For code, prove → test → smallest fix → CI → deploy → prove again.","Ask the user only for genuine account-level or break-glass actions."],1.0,1.45,11.3,4.75,18,WHITE,CYAN)
    txt(s,"LOOK → THINK → PRACTICE → PROVE",1.25,6.25,10.8,.55,24,GREEN,True,PP_ALIGN.CENTER)

    s = base_slide(prs, "Authority and scope")
    txt(s,"Primary current authorities",.75,1.42,5.2,.4,17,WHITE,True)
    bullets(s,["AGENTS.md","System3 Agent Runbook","Temporal Truth policy","Autonomous Operations policy","Master Goal Lock","Prediction Benchmark policy"],.8,1.95,5.4,3.85,13,MUTED,CYAN)
    txt(s,"This package",6.65,1.42,5.2,.4,17,WHITE,True)
    bullets(s,["Teaches safe use; it does not certify live health.","Covers all 22 canonical tabs and major repo workflows.","Generated from repository sources on 2026-08-27.","Revalidate URLs, commands and runtime claims when used.","CSV row IDs provide the detailed checklist."],6.7,1.95,5.65,3.85,13,MUTED,GREEN)
    txt(s,"Repository authority: psw2025-cmd/Genesis_System3 • Runtime: GCP system3-openalgo-safe / asia-south1 / genesis-system3-web • Broker: Dhan",.85,6.25,11.65,.5,11,AMBER,True,PP_ALIGN.CENTER)

    OUT.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_PATH)


def main():
    OUT.mkdir(parents=True, exist_ok=True)
    rows = build_rows()
    with CSV_PATH.open("w", newline="", encoding="utf-8-sig") as f:
        writer = csv.DictWriter(f, fieldnames=COLUMNS)
        writer.writeheader(); writer.writerows(rows)
    build_pptx()
    print(f"CSV={CSV_PATH} rows={len(rows)}")
    print(f"PPTX={PPTX_PATH}")


if __name__ == "__main__":
    main()
