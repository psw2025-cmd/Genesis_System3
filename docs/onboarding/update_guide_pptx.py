"""Idempotently append the ephemeral-storage warning slide."""

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
DECK = HERE / "SYSTEM3_NEW_USER_KID_LEVEL_GUIDE.pptx"
TITLE = "WARNING: Ephemeral Cloud Storage & Data Wipes"


def text(slide, value, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear(); frame.word_wrap = True
    paragraph = frame.paragraphs[0]; paragraph.alignment = align
    run = paragraph.add_run(); run.text = value
    run.font.name = "Aptos"; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color


def card(slide, x, title, body, accent):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.65), Inches(3.75), Inches(3.9))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(15, 31, 49); shape.line.color.rgb = accent
    text(slide, title, x + .25, 1.9, 3.25, .35, 15, RGBColor(243, 247, 252), True)
    text(slide, body, x + .25, 2.55, 3.25, 2.55, 13, RGBColor(167, 184, 204))


def main():
    prs = Presentation(DECK)
    titles = {shape.text.strip() for slide in prs.slides for shape in slide.shapes if getattr(shape, "has_text_frame", False)}
    if TITLE in titles:
        print("UNCHANGED: warning slide already exists")
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(7, 18, 31)
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(.1))
    banner.fill.solid(); banner.fill.fore_color.rgb = RGBColor(255, 73, 100); banner.line.color.rgb = RGBColor(255, 73, 100)
    text(slide, "SYSTEM3 • STORAGE SAFETY", .6, .35, 5.5, .25, 10, RGBColor(245, 165, 36), True)
    text(slide, TITLE, .6, .72, 12, .75, 26, RGBColor(243, 247, 252), True)
    card(slide, .7, "WHAT HAPPENS", "Cloud Run's normal container disk is temporary. A restart, new instance, or deployment can remove CSV/JSON files written inside the container.", RGBColor(255, 73, 100))
    card(slide, 4.78, "WHAT SURVIVES", "Only data explicitly written to a durable service—such as the approved Firestore backend—survives. A file path named state/ is not automatically durable.", RGBColor(32, 211, 238))
    card(slide, 8.87, "WHAT TO CHECK", "Trace both writer and reader. Confirm they use the same durable key/schema. Show EMPTY or BLOCKED honestly when lineage is missing; never manufacture market or ML proof.", RGBColor(24, 215, 130))
    text(slide, "Kid rule: a file inside the cloud box is like writing on a foggy window—it can disappear. Put important proof in the durable notebook.", .95, 6.05, 11.45, .7, 16, RGBColor(245, 165, 36), True, PP_ALIGN.CENTER)
    text(slide, "Educational warning • no LIVE/order permission", .65, 7.08, 6, .2, 8, RGBColor(167, 184, 204))
    text(slide, "LIVE OFF • PAPER / ANALYZE", 10, 7.05, 2.7, .25, 9, RGBColor(24, 215, 130), True, PP_ALIGN.RIGHT)
    prs.save(DECK)
    print(f"UPDATED: {DECK}")


if __name__ == "__main__":
    main()
